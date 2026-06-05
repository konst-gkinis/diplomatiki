# -*- coding: utf-8 -*-
"""
Generator for the thesis-defense presentation:
  «Πλατφόρμα λογισμικού για την επίτευξη του στόχου μηδενικές εκπομπές CO₂:
   Η περίπτωση της Mærsk»:  Κωνσταντίνος Γκίνης, Πανεπιστήμιο Πατρών.

Builds presentations/defense.pptx with python-pptx.

Design goals (per defense feedback):
  * Greek, University-of-Patras academic styling (burgundy + gold, serif headings).
  * Results-forward; theory compressed to a single slide.
  * Dedicated functional / non-functional (technical) requirements slides.
  * Screenshot PLACEHOLDERS mapped to the real platform features, swap in the
    actual images later by replacing add_placeholder(...) with add_image(...).
"""

import os
import unicodedata
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
SHOTS = os.path.join(HERE, "..", "screenshots")   # real platform screenshots
OUT = os.path.join(HERE, "defense.pptx")

# ── Palette (University of Patras) ────────────────────────────────────────
UP_RED   = RGBColor(0x7A, 0x16, 0x26)   # burgundy from the UP logo
UP_RED_D = RGBColor(0x57, 0x0F, 0x1B)   # darker burgundy
UP_GOLD  = RGBColor(0xC2, 0x9B, 0x52)   # gold/tan from the UP logo
DARK     = RGBColor(0x23, 0x23, 0x23)
GRAY     = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT    = RGBColor(0xF6, 0xF2, 0xEA)   # warm off-white
PANEL    = RGBColor(0xEC, 0xE4, 0xD6)   # warm panel
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PLACE_BG = RGBColor(0xF0, 0xED, 0xE6)
PLACE_BD = RGBColor(0xB9, 0xA5, 0x7E)

# Fonts bundled with LibreOffice (Greek-capable, metric-compatible with MS):
HEAD = "Lato"        # clean geometric sans for headings (native in Google Slides)
BODY = "Carlito"     # ≈ Calibri  (clean body text)
CODE = "DejaVu Sans Mono"   # monospace with Greek glyphs (for code snippets)

CODE_BG = RGBColor(0x27, 0x1C, 0x2F)   # dark eggplant (echoes Elixir purple)
CODE_FG = RGBColor(0xEC, 0xE6, 0xF0)
CODE_MUTE = RGBColor(0x9D, 0x90, 0xAB)  # comments

SW, SH = Inches(13.333), Inches(7.5)   # 16:9


def gr_upper(s):
    """Greek-aware uppercase: drop the monotonic tonos (and other accents),
    but keep the dialytika, per Greek orthography for all-caps text."""
    drop = {"̀", "́", "͂", "̀", "́", "̓"}
    decomposed = unicodedata.normalize("NFD", s.upper())
    kept = "".join(ch for ch in decomposed if ch not in drop)
    return unicodedata.normalize("NFC", kept)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── low-level helpers ─────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def _set_run(r, text, size, color, bold, italic, font):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    # ensure Greek glyphs use the same face for complex/east-asian slots
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs", "a:ea"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", font)


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            wrap=True):
    """lines: list of dicts or list of [runs]; each line -> paragraph.
    A 'line' may be:
      {"t":str,"size":..,"color":..,"bold":..,"italic":..,"font":..,
       "align":..,"bullet":bool,"space_after":..,"level":int}
    or a list of run-dicts (mixed formatting in one paragraph)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        runs = ln if isinstance(ln, list) else [ln]
        meta = runs[0]
        p.alignment = meta.get("align", align)
        if meta.get("space_after") is not None:
            p.space_after = Pt(meta["space_after"])
        if meta.get("space_before") is not None:
            p.space_before = Pt(meta["space_before"])
        if meta.get("line_spacing"):
            p.line_spacing = meta["line_spacing"]
        p.level = meta.get("level", 0)
        for rd in runs:
            r = p.add_run()
            _set_run(r, rd.get("t", ""), rd.get("size", 16),
                     rd.get("color", DARK), rd.get("bold", False),
                     rd.get("italic", False), rd.get("font", BODY))
        if meta.get("bullet"):
            _bullet(p, meta.get("bullet_color", UP_RED))
        else:
            _no_bullet(p)
    return tb


def _bullet(p, color):
    pPr = p._p.get_or_add_pPr()
    pPr.set("indent", str(-Inches(0.22)))
    pPr.set("marL", str(Inches(0.22)))
    buFont = pPr.makeelement(qn("a:buFont"),
                             {"typeface": "Arial", "pitchFamily": "34", "charset": "0"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
    pPr.append(buFont); pPr.append(buChar)


def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def fit(img_path, max_w, max_h):
    iw, ih = Image.open(img_path).size
    r = min(max_w / iw, max_h / ih)
    return Emu(int(iw * r)), Emu(int(ih * r))


def add_image(s, path, x, y, max_w, max_h, center_in=True):
    w, h = fit(path, max_w, max_h)
    if center_in:
        x = x + Emu(int((max_w - w) / 2))
        y = y + Emu(int((max_h - h) / 2))
    return s.shapes.add_picture(path, x, y, w, h)


def add_placeholder(s, x, y, w, h, title, desc):
    box = rect(s, x, y, w, h, fill=PLACE_BG, line=PLACE_BD, line_w=Pt(1.5))
    box.line.dash_style = None
    # dashed border
    ln = box.line._get_or_add_ln()
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"}); ln.append(d)
    textbox(s, x + Inches(0.2), y, w - Inches(0.4), h, [
        {"t": "▣  ΣΤΙΓΜΙΟΤΥΠΟ ΟΘΟΝΗΣ", "size": 13, "color": UP_RED,
         "bold": True, "font": BODY, "align": PP_ALIGN.CENTER, "space_after": 6},
        {"t": title, "size": 15, "color": DARK, "bold": True, "font": HEAD,
         "align": PP_ALIGN.CENTER, "space_after": 6},
        {"t": desc, "size": 11.5, "color": GRAY, "italic": True, "font": BODY,
         "align": PP_ALIGN.CENTER, "line_spacing": 1.05},
    ], anchor=MSO_ANCHOR.MIDDLE)


def _set_radius(shape, w, h, radius_in):
    """Pin a rounded-rectangle's corner radius to a fixed length (in inches),
    independent of the shape's size, so rounding stays subtle and consistent."""
    try:
        shape.adjustments[0] = min(0.5, float(Inches(radius_in)) /
                                   float(min(int(w), int(h))))
    except Exception:
        pass


def _add_shadow(shape, blur=Pt(3), dist=Pt(2), alpha=34000):
    """Attach an explicit soft outer drop shadow (renders in PowerPoint and
    Google Slides, not only in the LibreOffice PDF export)."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur)), "dist": str(int(dist)),
        "dir": "5400000", "rotWithShape": "0"})
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    shdw.append(clr); eff.append(shdw); spPr.append(eff)


def add_screenshot(s, path, x, y, w, h, caption=None, zoom=False):
    """Real platform screenshot in a framed 'window' with a burgundy title bar.

    Draw order matters: white card first, then the image, then the outline and
    caption bar ON TOP, so the border always frames the image cleanly and the
    image corners never overhang the frame.

    With zoom=True the shot is registered so a full-screen appendix slide is
    built later and the thumbnail becomes click-to-enlarge.
    """
    bar_h = Inches(0.34) if caption else Inches(0)
    pad = Inches(0.1)
    # 1) white background card, gently rounded, with a soft drop shadow
    card = rect(s, x, y, w, h, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_radius(card, w, h, 0.12)
    _add_shadow(card)
    # 2) the screenshot, inset and centred, behind the frame outline
    avail_w = w - 2 * pad
    avail_h = h - bar_h - 2 * pad
    iw, ih = fit(path, avail_w, avail_h)
    ix = x + Emu(int((w - iw) / 2))
    iy = y + bar_h + pad + Emu(int((avail_h - ih) / 2))
    pic = s.shapes.add_picture(path, ix, iy, iw, ih)
    # 3) outline drawn on top of the image (no fill, so clicks reach the image)
    out = rect(s, x, y, w, h, fill=None, line=PLACE_BD, line_w=Pt(1.0),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _set_radius(out, w, h, 0.12)
    # 4) caption bar on the very top (rounded to match, top corners only)
    bar = None
    if caption:
        bar = rect(s, x, y, w, bar_h, fill=UP_RED,
                   shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        _set_radius(bar, w, bar_h, 0.12)
        textbox(s, x + Inches(0.18), y, w - Inches(0.36), bar_h,
                [{"t": caption, "size": 11, "color": WHITE, "bold": True,
                  "font": BODY}], anchor=MSO_ANCHOR.MIDDLE)
    if zoom:
        ZOOMS.append({"pic": pic, "bar": bar, "src": s, "path": path,
                      "caption": caption or ""})
    return pic


def code_block(s, x, y, w, h, lines, fs=10.5):
    rect(s, x, y, w, h, fill=CODE_BG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.12),
                              w - Inches(0.4), h - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(1)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        _no_bullet(p)
        is_comment = ln.lstrip().startswith("#")
        r = p.add_run()
        _set_run(r, ln if ln else " ", fs,
                 CODE_MUTE if is_comment else CODE_FG, False, is_comment, CODE)
    return tb


# ── chrome: title bar + footer ────────────────────────────────────────────
PAGE = [0]
ZOOMS = []   # screenshots registered for click-to-enlarge appendix slides


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.12), fill=UP_RED)
    rect(s, 0, Inches(1.12), SW, Inches(0.07), fill=UP_GOLD)
    if kicker:
        textbox(s, Inches(0.55), Inches(0.13), Inches(11), Inches(0.3),
                [{"t": gr_upper(kicker), "size": 11, "color": UP_GOLD, "bold": True,
                  "font": BODY}])
    textbox(s, Inches(0.55), Inches(0.40), Inches(12.2), Inches(0.66),
            [{"t": title, "size": 25, "color": WHITE, "bold": True, "font": HEAD}],
            anchor=MSO_ANCHOR.MIDDLE)


def footer(s, dark=False):
    PAGE[0] += 1
    c = WHITE if dark else GRAY
    logo = "logo-ceid-white.png" if dark else "logo-ceid.png"
    add_image(s, os.path.join(ASSETS, logo), Inches(0.55), Inches(7.02),
              Inches(1.5), Inches(0.36), center_in=False)
    textbox(s, Inches(2.2), Inches(7.06), Inches(8.5), Inches(0.34),
            [{"t": "Πάγκος Εργασίας Εκπομπών · Διπλωματική Εργασία, Κ. Γκίνης",
              "size": 9.5, "color": c, "font": BODY}],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(12.2), Inches(7.06), Inches(0.9), Inches(0.34),
            [{"t": str(PAGE[0]), "size": 9.5, "color": c, "font": BODY,
              "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)


def body_area():
    """Return (x, y, w, h) of the usable content region below the header."""
    return Inches(0.55), Inches(1.42), Inches(12.23), Inches(5.5)


def chip(s, x, y, w, text, color=UP_RED):
    h = Inches(0.42)
    rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x, y, w, h, [{"t": text, "size": 13, "color": WHITE, "bold": True,
            "font": BODY, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)


# ── styled table ──────────────────────────────────────────────────────────
def table(s, x, y, w, rows, col_w, header_row=True, font_size=12.5,
          row_h=Inches(0.42)):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, x, y, w, row_h * nr).table
    gt.first_row = False
    gt.horz_banding = False
    # remove default style banding by setting explicit fills
    total = sum(col_w)
    for j, cw in enumerate(col_w):
        gt.columns[j].width = Emu(int(w * cw / total))
    for i, row in enumerate(rows):
        gt.rows[i].height = row_h
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if header_row and i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = UP_RED
                col, bold, fs = WHITE, True, font_size
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = WHITE if (i % 2) else PANEL
                col, bold, fs = DARK, False, font_size
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.LEFT
            r = p.add_run()
            _set_run(r, str(val), fs, col, bold or (j == 0 and i > 0 and val.startswith("•")),
                     False, BODY)
            if j == 0 and i > 0:
                r.font.bold = True
                r.font.color.rgb = UP_RED_D
    return gt


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1, TITLE
# ══════════════════════════════════════════════════════════════════════════
def s_title():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=LIGHT)
    rect(s, 0, 0, Inches(0.32), SH, fill=UP_RED)
    rect(s, Inches(0.32), 0, Inches(0.07), SH, fill=UP_GOLD)
    add_image(s, os.path.join(ASSETS, "logo-up-emblem.png"),
              Inches(0.9), Inches(0.55), Inches(1.7), Inches(1.7), center_in=False)
    add_image(s, os.path.join(ASSETS, "logo-ceid.png"),
              Inches(9.55), Inches(0.95), Inches(2.85), Inches(0.95),
              center_in=False)
    textbox(s, Inches(2.85), Inches(0.62), Inches(6.4), Inches(1.4), [
        {"t": "ΠΑΝΕΠΙΣΤΗΜΙΟ ΠΑΤΡΩΝ", "size": 17, "color": UP_RED, "bold": True,
         "font": HEAD, "space_after": 2},
        {"t": "Πολυτεχνική Σχολή", "size": 13, "color": DARK, "font": BODY,
         "space_after": 1},
        {"t": "Τμήμα Μηχανικών Η/Υ και Πληροφορικής", "size": 13, "color": DARK,
         "font": BODY},
    ])
    rect(s, Inches(0.9), Inches(2.62), Inches(11.5), Pt(2), fill=UP_GOLD)
    textbox(s, Inches(0.9), Inches(2.95), Inches(11.6), Inches(2.0), [
        {"t": "Πλατφόρμα λογισμικού για την επίτευξη του στόχου",
         "size": 30, "color": UP_RED_D, "bold": True, "font": HEAD,
         "line_spacing": 1.05, "space_after": 0},
        {"t": "μηδενικές εκπομπές CO₂: Η περίπτωση της Mærsk",
         "size": 30, "color": UP_RED_D, "bold": True, "font": HEAD,
         "line_spacing": 1.05},
    ])
    textbox(s, Inches(0.9), Inches(4.55), Inches(11.6), Inches(0.5), [
        {"t": "Διπλωματική Εργασία", "size": 15, "color": GRAY, "italic": True,
         "font": BODY}])
    # author + committee block
    textbox(s, Inches(0.9), Inches(5.25), Inches(6.5), Inches(1.6), [
        {"t": "Κωνσταντίνος Γκίνης", "size": 19, "color": DARK, "bold": True,
         "font": HEAD, "space_after": 8},
        [{"t": "Επιβλέπων:  ", "size": 13, "color": GRAY, "font": BODY},
         {"t": "Χρήστος Μπούρας, Καθηγητής", "size": 13, "color": DARK, "font": BODY}],
    ])
    textbox(s, Inches(7.4), Inches(5.25), Inches(5.0), Inches(1.6), [
        {"t": "Εξεταστική Επιτροπή", "size": 12, "color": GRAY, "bold": True,
         "font": BODY, "space_after": 4},
        {"t": "Ιωάννης Γαροφαλάκης, Καθηγητής", "size": 12, "color": DARK,
         "font": BODY, "space_after": 1},
        {"t": "Εύη Παπαϊωάννου, Επίκουρη Καθηγήτρια", "size": 12, "color": DARK,
         "font": BODY},
    ])
    textbox(s, Inches(0.9), Inches(6.85), Inches(11.6), Inches(0.4), [
        {"t": "Πάτρα, Ιούνιος 2026", "size": 12, "color": UP_RED, "bold": True,
         "font": BODY}])


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2, THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
def s_problem():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Το πλαίσιο", "Το πρόβλημα: η χειρωνακτική διαδικασία")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(6.7), Inches(0.9), [
        {"t": "Η Maersk (στόλος 700+ πλοίων) αδυνατούσε να ανταποκριθεί "
              "αξιόπιστα στις απαιτήσεις για δεδομένα εκπομπών.",
         "size": 16, "color": DARK, "font": BODY, "line_spacing": 1.12}])
    steps = [
        ("Αίτημα μέσω email", "Ο πελάτης ή το τμήμα πωλήσεων ζητά δεδομένα εκπομπών"),
        ("Χειρωνακτική εξαγωγή", "Μη αυτόματη άντληση δεδομένων από πολλαπλές πηγές"),
        ("Excel + συντελεστές", "Εφαρμογή συντελεστών εκπομπών σε υπολογιστικά φύλλα"),
        ("Απάντηση σε 3–8 εβδομάδες", "Αποστολή αποτελέσματος ξανά μέσω email"),
    ]
    yy = by + Inches(1.0)
    for i, (t, d) in enumerate(steps):
        rect(s, bx, yy, Inches(0.42), Inches(0.42), fill=UP_RED,
             shape=MSO_SHAPE.OVAL)
        textbox(s, bx, yy, Inches(0.42), Inches(0.42),
                [{"t": str(i + 1), "size": 14, "color": WHITE, "bold": True,
                  "font": BODY, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, bx + Inches(0.6), yy - Inches(0.04), Inches(6.0), Inches(0.6), [
            [{"t": t + ":  ", "size": 14.5, "color": UP_RED_D, "bold": True,
              "font": BODY},
             {"t": d, "size": 13, "color": DARK, "font": BODY}]])
        yy = yy + Inches(0.72)
    # right panel: the cost
    px = Inches(7.7)
    rect(s, px, by, Inches(5.05), Inches(4.55), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, px, by, Inches(0.12), Inches(4.55), fill=UP_GOLD)
    textbox(s, px + Inches(0.4), by + Inches(0.3), Inches(4.3), Inches(4.0), [
        {"t": "Τι κόστιζε αυτό", "size": 17, "color": UP_RED, "bold": True,
         "font": HEAD, "space_after": 12},
        {"t": "Ακρίβεια", "bullet": True, "size": 14.5, "color": DARK, "bold": True,
         "font": BODY, "space_after": 1},
        {"t": "λάθη μεθοδολογίας, ελλιπής κανονικοποίηση δεδομένων", "size": 12.5,
         "color": GRAY, "font": BODY, "level": 1, "space_after": 8},
        {"t": "Ταχύτητα", "bullet": True, "size": 14.5, "color": DARK, "bold": True,
         "font": BODY, "space_after": 1},
        {"t": "καθυστέρηση εβδομάδων, απώλεια συμβολαίων", "size": 12.5,
         "color": GRAY, "font": BODY, "level": 1, "space_after": 8},
        {"t": "Ιχνηλασιμότητα", "bullet": True, "size": 14.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "καμία δυνατότητα ανεξάρτητου ελέγχου (audit)", "size": 12.5,
         "color": GRAY, "font": BODY, "level": 1, "space_after": 8},
        {"t": "Κλίμακα", "bullet": True, "size": 14.5, "color": DARK, "bold": True,
         "font": BODY, "space_after": 1},
        {"t": "1.300+ χρήστες, ανθρώπινη δυναμικότητα δεν κλιμακώνεται", "size": 12.5,
         "color": GRAY, "font": BODY, "level": 1},
    ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3, GIA AUDIT (catalyst → requirements)
# ══════════════════════════════════════════════════════════════════════════
def s_audit():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Ο καταλύτης", "Ο εσωτερικός έλεγχος GIA 2023")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(12.2), Inches(0.7), [
        [{"t": "Ο έλεγχος Group Internal Audit (2023) ", "size": 15.5,
          "color": DARK, "font": BODY},
         {"t": "μετέτρεψε ένα γνωστό λειτουργικό πρόβλημα σε επίσημη, "
               "επείγουσα προτεραιότητα, και όρισε τις απαιτήσεις του συστήματος.",
          "size": 15.5, "color": DARK, "font": BODY}]],
    )
    cards = [
        ("Παρακολούθηση δεδομένων", "Καμία ενιαία ιχνηλάτηση ροής δεδομένων από την πηγή έως την αναφορά"),
        ("Ίχνος ελέγχου (audit trail)", "Αδύνατη η ανεξάρτητη επαλήθευση των υπολογισμών, κρίσιμο για CSRD/EU ETS"),
        ("Κλιμακωσιμότητα", "Αύξηση αιτημάτων χωρίς αντίστοιχη κλιμάκωση ανθρώπινης δυναμικότητας"),
        ("Ποσοστό σφαλμάτων", "Εγγενώς υψηλότερα λάθη σε χειρωνακτικές διαδικασίες"),
    ]
    cw, ch = Inches(2.92), Inches(2.55)
    gap = Inches(0.16)
    x0 = bx
    yy = by + Inches(1.0)
    for i, (t, d) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, yy, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, yy, cw, Inches(0.1), fill=UP_RED)
        textbox(s, x + Inches(0.22), yy + Inches(0.28), cw - Inches(0.44),
                ch - Inches(0.4), [
            {"t": "0" + str(i + 1), "size": 22, "color": UP_GOLD, "bold": True,
             "font": HEAD, "space_after": 6},
            {"t": t, "size": 15, "color": UP_RED_D, "bold": True, "font": HEAD,
             "space_after": 8, "line_spacing": 1.0},
            {"t": d, "size": 12.5, "color": DARK, "font": BODY, "line_spacing": 1.1},
        ])
    strip = rect(s, bx, yy + ch + Inches(0.22), Inches(12.2), Inches(0.6), fill=PANEL,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_shadow(strip)
    textbox(s, bx + Inches(0.3), yy + ch + Inches(0.22), Inches(11.6), Inches(0.6), [
        [{"t": "➜  Αποτέλεσμα:  ", "size": 14, "color": UP_RED, "bold": True,
          "font": BODY},
         {"t": "έγκριση πόρων για εξ ολοκλήρου νέο σύστημα, ανάπτυξη από τον "
               "Μάρτιο 2023, πρώτοι χρήστες παραγωγής τον Σεπτέμβριο 2023.",
          "size": 14, "color": DARK, "font": BODY}]],
        anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4, FUNCTIONAL REQUIREMENTS
# ══════════════════════════════════════════════════════════════════════════
def s_func_req():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Προδιαγραφές · 1/2", "Λειτουργικές απαιτήσεις")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(12.2), Inches(0.55), [
        [{"t": "Από τις ανάγκες ", "size": 14.5, "color": DARK, "font": BODY},
         {"t": "έξι ομάδων χρηστών", "size": 14.5, "color": UP_RED_D, "bold": True,
          "font": BODY},
         {"t": " προκύπτουν οι λειτουργίες που πρέπει να προσφέρει το σύστημα:",
          "size": 14.5, "color": DARK, "font": BODY}]])
    # ── hero: the derived functional capabilities (numbered) ──
    funcs = [
        ("Αναζήτηση εκπομπών", "ανά πελάτη και χρονικό εύρος, σε δευτερόλεπτα"),
        ("Ανάλυση ανά αποστολή / δρομολόγιο", "ποια trade lanes βαραίνουν το αποτύπωμα"),
        ("Εξαγωγή δεδομένων σε Excel", "η μορφή που χρησιμοποιούν στην πράξη οι χρήστες"),
        ("Αυτόματη έκδοση πιστοποιητικών", "έγγραφα ECO σε PDF, χωρίς χειρωνακτική εργασία"),
        ("Self-service", "χωρίς μεσολάβηση ειδικής ομάδας αναφορών"),
    ]
    yy = by + Inches(0.7)
    for i, (t, d) in enumerate(funcs):
        rect(s, bx, yy, Inches(0.46), Inches(0.46), fill=UP_RED, shape=MSO_SHAPE.OVAL)
        textbox(s, bx, yy, Inches(0.46), Inches(0.46),
                [{"t": str(i + 1), "size": 15, "color": WHITE, "bold": True,
                  "font": BODY, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, bx + Inches(0.66), yy - Inches(0.03), Inches(7.4), Inches(0.6), [
            [{"t": t + "  ", "size": 16, "color": UP_RED_D, "bold": True,
              "font": HEAD},
             {"t": d, "size": 12.5, "color": GRAY, "font": BODY}]],
            anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.78)
    # ── compact: the six user groups (de-emphasised) ──
    px = Inches(8.6)
    rect(s, px, by + Inches(0.7), Inches(4.2), Inches(3.9), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, px, by + Inches(0.7), Inches(0.12), Inches(3.9), fill=UP_GOLD)
    groups = [
        ("ECO Delivery Products", "πιστοποιητικά ανά αποστολή"),
        ("Regional Product Mgmt", "τάσεις ανά περιοχή"),
        ("Commercial Sustainability", "στοιχεία ESG / CSRD"),
        ("Regional Sales", "δεδομένα σε διαπραγμάτευση"),
        ("Contract Management", "audit-ready αποδείξεις"),
        ("Account Managers", "εκπομπές πελάτη real-time"),
    ]
    gi = [{"t": "Οι έξι ομάδες χρηστών", "size": 14, "color": UP_RED, "bold": True,
           "font": HEAD, "space_after": 9}]
    for t, d in groups:
        gi.append([{"t": t, "size": 12, "color": UP_RED_D, "bold": True,
                    "font": BODY},
                   {"t": "  " + d, "size": 11, "color": GRAY, "font": BODY}])
        gi[-1][0]["bullet"] = True
        gi[-1][0]["space_after"] = 8
        gi[-1][0]["line_spacing"] = 1.0
    textbox(s, px + Inches(0.34), by + Inches(0.92), Inches(3.7), Inches(3.5), gi)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5, NON-FUNCTIONAL / TECHNICAL REQUIREMENTS
# ══════════════════════════════════════════════════════════════════════════
def s_nonfunc_req():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Προδιαγραφές · 2/2", "Μη-λειτουργικές & τεχνικές απαιτήσεις")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(12.2), Inches(0.5), [
        {"t": "Οι ανάγκες των χρηστών και τα ευρήματα GIA ορίζουν πώς πρέπει να "
              "συμπεριφέρεται το σύστημα:", "size": 14.5, "color": DARK, "font": BODY}])
    items = [
        ("Ιχνηλασιμότητα", "Πλήρης διαδρομή δεδομένων από την πηγή έως κάθε αναφορά"),
        ("Αμετάβλητο audit trail", "Καταγραφή κάθε αλλαγής, ανεξάρτητα επαληθεύσιμη"),
        ("Δεδομένα πραγματικού χρόνου", "Freshness σε ώρες αντί για εβδομάδες"),
        ("Κλιμακωσιμότητα", "1.300+ χρήστες, αυξανόμενος όγκος αιτημάτων"),
        ("Υψηλή διαθεσιμότητα", "Production-grade, χωρίς διακοπές υπηρεσίας"),
        ("Επαληθευσιμότητα & συμμόρφωση", "Έτοιμο για CSRD, EU ETS, ISCC, CDP"),
        ("Αυτοματοποίηση", "Πιστοποιητικά χωρίς ανθρώπινη παρέμβαση"),
        ("Ασφάλεια & εξουσιοδότηση", "Πρόσβαση μόνο με εταιρικά δικαιώματα"),
    ]
    cw, ch = Inches(2.95), Inches(1.32)
    gx, gy = Inches(0.16), Inches(0.16)
    x0, y0 = bx, by + Inches(0.6)
    for i, (t, d) in enumerate(items):
        col, row = i % 4, i // 4
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.1), ch, fill=UP_RED)
        textbox(s, x + Inches(0.26), y + Inches(0.14), cw - Inches(0.4),
                ch - Inches(0.2), [
            {"t": t, "size": 14, "color": UP_RED_D, "bold": True, "font": HEAD,
             "space_after": 4, "line_spacing": 1.0},
            {"t": d, "size": 11.8, "color": DARK, "font": BODY, "line_spacing": 1.05},
        ])
    # bottom strip: how the design answers them
    yb = y0 + 2 * (ch + gy) + Inches(0.08)
    strip = rect(s, bx, yb, Inches(12.2), Inches(0.62), fill=PANEL,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_shadow(strip)
    textbox(s, bx + Inches(0.3), yb, Inches(11.6), Inches(0.62), [
        [{"t": "Τεχνική απάντηση:  ", "size": 13.5, "color": UP_RED, "bold": True,
          "font": BODY},
         {"t": "event-driven αρχιτεκτονική · event sourcing για το audit trail · "
               "Phoenix LiveView για real-time · BEAM/Elixir για διαθεσιμότητα.",
          "size": 13.5, "color": DARK, "font": BODY}]],
        anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6, THE SOLUTION: 4 PILLARS
# ══════════════════════════════════════════════════════════════════════════
def s_pillars():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Η λύση", "Emissions Workbench: τέσσερις πυλώνες")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(12.2), Inches(0.5), [
        [{"t": "Μία πλατφόρμα ", "size": 15, "color": DARK, "font": BODY},
         {"t": "(emissions-workbench.maersk.com)", "size": 13, "color": GRAY,
          "italic": True, "font": BODY},
         {"t": ", σε ένα monorepo, με κοινή ροή δεδομένων", "size": 15,
          "color": DARK, "font": BODY}]])
    pillars = [
        ("Ocean Emissions", "Αυτόματη μέτρηση εκπομπών ανά αποστολή από τηλεμετρία στόλου (STAR Connect), με EU ETS"),
        ("ECO Delivery", "Εμπορικά προϊόντα χαμηλού άνθρακα, αυτόματα πιστοποιητικά, Energy Bank (mass balance)"),
        ("Net Zero 2040", "Εσωτερικοί πίνακες για οδικούς χάρτες αποκαρβονοποίησης & δεσμεύσεις SBTi"),
        ("Bunker Optimization", "Βελτιστοποίηση πλάνων ανεφοδιασμού για ελάχιστο κόστος καυσίμου σε περιβάλλον ETS"),
    ]
    cw, ch = Inches(2.92), Inches(2.85)
    gap = Inches(0.16)
    yy = by + Inches(0.65)
    for i, (t, d) in enumerate(pillars):
        x = bx + i * (cw + gap)
        rect(s, x, yy, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, yy, cw, Inches(0.8), fill=UP_RED,
             shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        textbox(s, x + Inches(0.18), yy, cw - Inches(0.36), Inches(0.8),
                [{"t": t, "size": 15, "color": WHITE, "bold": True, "font": HEAD,
                  "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(0.24), yy + Inches(0.98), cw - Inches(0.48),
                Inches(1.7), [
            {"t": d, "size": 13, "color": DARK, "font": BODY, "line_spacing": 1.14}])
    # ── data-flow mini graphic ──
    fy = yy + ch + Inches(0.28)
    textbox(s, bx, fy - Inches(0.04), Inches(4), Inches(0.3),
            [{"t": "ΡΟΗ ΔΕΔΟΜΕΝΩΝ", "size": 10.5, "color": UP_RED, "bold": True,
              "font": BODY}])

    def fnode(x, y, w, text, fill, tcol, fs=11.5):
        nd = rect(s, x, y, w, Inches(0.5), fill=fill,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_shadow(nd)
        textbox(s, x + Inches(0.05), y, w - Inches(0.1), Inches(0.5),
                [{"t": text, "size": fs, "color": tcol, "bold": True, "font": BODY,
                  "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
                anchor=MSO_ANCHOR.MIDDLE)

    def farrow(x, y):
        a = rect(s, x, y, Inches(0.5), Inches(0.28), fill=UP_GOLD,
                 shape=MSO_SHAPE.RIGHT_ARROW)
        return a

    ny = fy + Inches(0.32)
    src_w = Inches(2.45)
    fnode(bx, ny - Inches(0.32), src_w, "Ocean Emissions", PANEL, UP_RED_D)
    fnode(bx, ny + Inches(0.32), src_w, "Bunker Optimization", PANEL, UP_RED_D)
    farrow(bx + src_w + Inches(0.12), ny + Inches(0.02))
    eb_x = bx + src_w + Inches(0.74)
    fnode(eb_x, ny, Inches(2.5), "Energy Bank", UP_RED, WHITE, fs=12.5)
    farrow(eb_x + Inches(2.62), ny + Inches(0.02))
    cert_x = eb_x + Inches(3.24)
    fnode(cert_x, ny, Inches(2.7), "Πιστοποιητικά ECO", UP_GOLD, RGBColor(0x3A, 0x2A, 0x05), fs=12.5)
    textbox(s, cert_x + Inches(2.85), ny - Inches(0.02), Inches(1.6), Inches(0.55),
            [{"t": "επικοινωνία", "size": 10, "color": GRAY, "font": BODY,
              "line_spacing": 1.0},
             {"t": "μέσω Kafka", "size": 10, "color": GRAY, "italic": True,
              "font": BODY}], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7, ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════
def s_arch():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Τεχνική επισκόπηση", "Αρχιτεκτονική του συστήματος")
    bx, by, bw, bh = body_area()
    add_image(s, os.path.join(ASSETS, "system_overview_diagram.png"),
              bx, by, Inches(8.5), Inches(5.0), center_in=True)
    px = Inches(9.2)
    rect(s, px, by, Inches(3.6), Inches(5.0), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, px, by, Inches(0.12), Inches(5.0), fill=UP_GOLD)
    textbox(s, px + Inches(0.34), by + Inches(0.3), Inches(3.1), Inches(4.5), [
        {"t": "Βασικές αρχές", "size": 16, "color": UP_RED, "bold": True,
         "font": HEAD, "space_after": 12},
        {"t": "Clean Architecture", "bullet": True, "size": 13.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "domain ανεξάρτητο από υποδομή", "size": 11.5, "color": GRAY,
         "font": BODY, "level": 1, "space_after": 9},
        {"t": "Event-driven ingestion", "bullet": True, "size": 13.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "Kafka + Oban workers", "size": 11.5, "color": GRAY, "font": BODY,
         "level": 1, "space_after": 9},
        {"t": "Event Sourcing (Energy Bank)", "bullet": True, "size": 13.5,
         "color": DARK, "bold": True, "font": BODY, "space_after": 1},
        {"t": "πλήρες ιστορικό & audit trail", "size": 11.5, "color": GRAY,
         "font": BODY, "level": 1, "space_after": 9},
        {"t": "Real-time web (LiveView)", "bullet": True, "size": 13.5,
         "color": DARK, "bold": True, "font": BODY, "space_after": 1},
        {"t": "ενημέρωση χωρίς reload", "size": 11.5, "color": GRAY, "font": BODY,
         "level": 1, "space_after": 9},
        {"t": "Azure / Kubernetes", "bullet": True, "size": 13.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "auto-scaling, IaC με Terraform", "size": 11.5, "color": GRAY,
         "font": BODY, "level": 1},
    ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDES 8–10, PLATFORM SCREENSHOTS (placeholders)
# ══════════════════════════════════════════════════════════════════════════
def _demo_takeaway(s, text):
    y = Inches(6.06)
    bar = rect(s, Inches(0.55), y, Inches(12.23), Inches(0.44), fill=PANEL,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_shadow(bar)
    textbox(s, Inches(0.85), y, Inches(11.7), Inches(0.44),
            [[{"t": "➜  ", "size": 12.5, "color": UP_RED, "bold": True, "font": BODY},
              {"t": text, "size": 12.5, "color": DARK, "font": BODY}]],
            anchor=MSO_ANCHOR.MIDDLE)


def s_demo1():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Η πλατφόρμα στην πράξη · 1/3", "Αναζήτηση & ανάλυση εκπομπών")
    bx, by, bw, bh = body_area()
    half = Inches(5.965)
    gap = Inches(0.3)
    sh_y = by + Inches(0.05)
    sh_h = Inches(4.5)
    add_screenshot(s, os.path.join(SHOTS, "1_search_results.png"), bx, sh_y,
                   half, sh_h, "Αναζήτηση εκπομπών πελάτη", zoom=True)
    add_screenshot(s, os.path.join(SHOTS, "2_routes.png"),
                   bx + half + gap, sh_y, half, sh_h,
                   "Ανάλυση ανά αποστολή / δρομολόγιο", zoom=True)
    _demo_takeaway(s, "Self-service αναζήτηση και drill-down ανά trade lane: "
                      "από email-και-αναμονή 3–8 εβδομάδων σε άμεση απάντηση.")
    footer(s)


def s_demo2():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Η πλατφόρμα στην πράξη · 2/3", "Πιστοποιητικά & εξαγωγή")
    bx, by, bw, bh = body_area()
    sh_y = by + Inches(0.05)
    left_w = Inches(6.1)
    add_screenshot(s, os.path.join(SHOTS, "Certificate.png"), bx, sh_y,
                   left_w, Inches(4.28), "Παραγόμενο πιστοποιητικό ECO (PDF)",
                   zoom=True)
    rx = bx + left_w + Inches(0.3)
    rw = Inches(5.83)
    add_screenshot(s, os.path.join(SHOTS, "3_oban_progress.png"), rx, sh_y, rw,
                   Inches(2.32), "Oban: ασύγχρονες εργασίες (παραγωγή & ingestion)",
                   zoom=True)
    add_screenshot(s, os.path.join(SHOTS, "5_excel.png"), rx,
                   sh_y + Inches(2.5), rw, Inches(1.78), "Εξαγωγή δεδομένων σε Excel",
                   zoom=True)
    _demo_takeaway(s, "Oban workers αντλούν από το Energy Bank, επαληθεύουν "
                      "υπόλοιπα και παράγουν έγγραφα ISCC: από εβδομάδες σε λεπτά.")
    footer(s)


def s_demo3():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Η πλατφόρμα στην πράξη · 3/3", "Energy Bank & Bunker Optimization")
    bx, by, bw, bh = body_area()
    half = Inches(5.965)
    gap = Inches(0.3)
    sh_y = by + Inches(0.05)
    sh_h = Inches(4.5)
    add_screenshot(s, os.path.join(SHOTS, "6_energy_bank_transactions.png"),
                   bx, sh_y, half, sh_h, "Energy Bank: event-sourced ledger",
                   zoom=True)
    add_screenshot(s, os.path.join(SHOTS, "7_bunker_plan.png"),
                   bx + half + gap, sh_y, half, sh_h,
                   "Bunker Optimization: πλάνο ανεφοδιασμού", zoom=True)
    _demo_takeaway(s, "Αμετάβλητο ιστορικό συναλλαγών (audit trail) στο ισοζύγιο "
                      "μάζας · βέλτιστο πλάνο ανεφοδιασμού από τον solver.")
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11, THEORY (compressed to one slide)
# ══════════════════════════════════════════════════════════════════════════
def s_theory():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Θεωρητικό υπόβαθρο (σύνοψη)", "Λογιστική εκπομπών: τα βασικά")
    bx, by, bw, bh = body_area()
    add_image(s, os.path.join(ASSETS, "Carbon_Accounting_Scopes.png"),
              bx, by + Inches(0.2), Inches(6.2), Inches(4.6), center_in=True)
    px = Inches(7.1)
    concepts = [
        ("Scope 1 / 2 / 3", "Άμεσες · έμμεσες από ενέργεια · αλυσίδα αξίας"),
        ("TTW vs WTW", "Όρια συστήματος: Tank-to-Wake vs Well-to-Wake"),
        ("Trade factors", "Συντελεστές εμπορικής διαδρομής, ο πυρήνας του υπολογισμού"),
        ("Mass balance", "Λογιστική σύνδεση βιώσιμου καυσίμου με αξιώσεις"),
        ("Book-and-claim", "Αναφορά αποσυνδεδεμένη από τη φυσική ροή"),
        ("ISCC / GLEC / GHG Protocol", "Τα πρότυπα πιστοποίησης & μεθοδολογίας"),
    ]
    items = [{"t": "Από τη θεωρία στον κώδικα", "size": 16, "color": UP_RED,
              "bold": True, "font": HEAD, "space_after": 12}]
    for t, d in concepts:
        items.append([{"t": t + ":  ", "size": 14, "color": UP_RED_D,
                       "bold": True, "font": BODY},
                      {"t": d, "size": 13, "color": DARK, "font": BODY}])
        items[-1][0]["bullet"] = True
        items[-1][0]["space_after"] = 9
        items[-1][0]["line_spacing"] = 1.05
    textbox(s, px, by + Inches(0.15), Inches(5.6), Inches(4.6), items)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 12, TECHNOLOGIES
# ══════════════════════════════════════════════════════════════════════════
def s_tech():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Τεχνική επισκόπηση", "Τεχνολογίες")
    bx, by, bw, bh = body_area()
    # ── hero row: the three flagship technologies, with logos ──
    heroes = [
        ("logo-elixir.png", "BEAM / Elixir", "Υψηλή διαθεσιμότητα & ταυτοχρονισμός"),
        ("logo-phoenix.png", "Phoenix LiveView", "Real-time web χωρίς reload σελίδας"),
        ("logo-nix.png", "Nix / Flakes", "Reproducible builds & dev environments"),
    ]
    hcw = Inches(3.92)
    hgx = Inches(0.23)
    hch = Inches(2.35)
    hy = by + Inches(0.05)
    for i, (logo, t, d) in enumerate(heroes):
        x = bx + i * (hcw + hgx)
        rect(s, x, hy, hcw, hch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, hy, hcw, Inches(0.1), fill=UP_GOLD)
        add_image(s, os.path.join(ASSETS, logo), x + Inches(0.4),
                  hy + Inches(0.32), hcw - Inches(0.8), Inches(0.95))
        textbox(s, x + Inches(0.2), hy + Inches(1.5), hcw - Inches(0.4),
                Inches(0.8), [
            {"t": t, "size": 17, "color": UP_RED_D, "bold": True, "font": HEAD,
             "align": PP_ALIGN.CENTER, "space_after": 3},
            {"t": d, "size": 12.5, "color": DARK, "font": BODY,
             "align": PP_ALIGN.CENTER, "line_spacing": 1.05},
        ])
    # ── supporting stack: compact cards ──
    techs = [
        ("PostgreSQL / Ecto", "Σχεσιακή μονιμότητα δεδομένων"),
        ("Apache Kafka", "Event-driven ροή μηνυμάτων"),
        ("GitHub Actions", "CI/CD: αυτόματα builds & deploys"),
        ("Azure · Kubernetes · Terraform", "Cloud, auto-scaling & IaC"),
    ]
    cw, ch = Inches(2.9), Inches(1.2)
    gx = Inches(0.21)
    y0 = hy + hch + Inches(0.32)
    for i, (t, d) in enumerate(techs):
        x = bx + i * (cw + gx)
        rect(s, x, y0, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y0, Inches(0.1), ch, fill=UP_GOLD)
        textbox(s, x + Inches(0.24), y0 + Inches(0.16), cw - Inches(0.4),
                ch - Inches(0.2), [
            {"t": t, "size": 13, "color": UP_RED_D, "bold": True, "font": HEAD,
             "space_after": 4, "line_spacing": 1.0},
            {"t": d, "size": 11.5, "color": DARK, "font": BODY, "line_spacing": 1.02},
        ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE, CODE HIGHLIGHTS (Elixir expressiveness)
# ══════════════════════════════════════════════════════════════════════════
def s_code():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Τεχνική επισκόπηση", "Εκφραστικότητα: Elixir στην πράξη")
    bx, by, bw, bh = body_area()
    lw = Inches(7.7)

    textbox(s, bx, by, lw, Inches(0.3), [
        {"t": "Pattern matching σε πολλαπλές ρήτρες: επιλογή & "
              "κανονικοποίηση μονάδων", "size": 12, "color": UP_RED, "bold": True,
         "font": BODY}])
    code_block(s, bx, by + Inches(0.32), lw, Inches(1.95), [
        "# Το pattern matching επιλέγει τη ρήτρα και αποδομεί τα δεδομένα",
        "def calculate(:diesel,",
        "      %{energy_ttw: %Measurement{unit: :megajoules} = e} = a),",
        "  do: calculate(:diesel, %{a | energy_ttw: convert(e, :kwh)})",
        "",
        "def calculate(:diesel,",
        "      %{energy_ttw: %Measurement{unit: :kwh}} = args) do",
        "  # ... ο καθαρός υπολογισμός CO₂e",
        "end",
    ], fs=10.5)

    textbox(s, bx, by + Inches(2.45), lw, Inches(0.3), [
        {"t": "Pipe operator |>  &  ερώτημα Ecto (πρόσβαση σε δεδομένα)",
         "size": 12, "color": UP_RED, "bold": True, "font": BODY}])
    code_block(s, bx, by + Inches(2.77), lw, Inches(2.55), [
        "# |> : ροή μετασχηματισμών, αριστερά προς δεξιά",
        "kwh",
        "|> Measurement.value!(:kwh)",
        "|> Decimal.div(ratio_kwh_to_litres(ratio, factors))",
        "|> Measurement.new(:litres)",
        "",
        "# Ecto: δηλωτικό, type-safe, συνθέσιμο ερώτημα",
        "from(f in RemainingFuelOnBoard,",
        "  where:  f.vessel_id == ^vessel_id,",
        "  select: max(f.updated_at))",
        "|> Repo.one()",
    ], fs=10.5)

    # right: why functional
    px = bx + lw + Inches(0.25)
    pw = Inches(4.28)
    rect(s, px, by + Inches(0.32), pw, Inches(5.0), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, px, by + Inches(0.32), Inches(0.12), Inches(5.0), fill=UP_GOLD)
    textbox(s, px + Inches(0.34), by + Inches(0.6), pw - Inches(0.6), Inches(4.5), [
        {"t": "Γιατί λειτουργικός;", "size": 16, "color": UP_RED, "bold": True,
         "font": HEAD, "space_after": 12},
        {"t": "Pattern matching", "bullet": True, "size": 13.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "επιλογή ρήτρας & αποδόμηση δεδομένων, δηλωτικά, χωρίς if",
         "size": 11.8, "color": GRAY, "font": BODY, "level": 1, "space_after": 10},
        {"t": "Pipe |>", "bullet": True, "size": 13.5, "color": DARK, "bold": True,
         "font": BODY, "space_after": 1},
        {"t": "αναγνώσιμη ροή, χωρίς ενδιάμεσες μεταβλητές", "size": 11.8,
         "color": GRAY, "font": BODY, "level": 1, "space_after": 10},
        {"t": "Αμεταβλητότητα", "bullet": True, "size": 13.5, "color": DARK,
         "bold": True, "font": BODY, "space_after": 1},
        {"t": "ασφάλεια στον ταυτοχρονισμό του BEAM", "size": 11.8, "color": GRAY,
         "font": BODY, "level": 1, "space_after": 10},
        {"t": "Ecto", "bullet": True, "size": 13.5, "color": DARK, "bold": True,
         "font": BODY, "space_after": 1},
        {"t": "συνθέσιμα, type-safe ερωτήματα προς PostgreSQL", "size": 11.8,
         "color": GRAY, "font": BODY, "level": 1},
    ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 13, WORKING METHODS (XP)
# ══════════════════════════════════════════════════════════════════════════
def s_methods():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Εργασιακές μέθοδοι", "Extreme Programming σε ομάδα 28 μελών")
    bx, by, bw, bh = body_area()
    add_image(s, os.path.join(ASSETS, "Extreme_Programming_Loops.png"),
              bx, by, Inches(5.2), Inches(5.0), center_in=True)
    px = Inches(6.0)
    practices = [
        ("Pair programming", "συνεχής έλεγχος κώδικα & μεταφορά γνώσης"),
        ("TDD", "ανάπτυξη καθοδηγούμενη από δοκιμές, 80%+ κάλυψη"),
        ("CI/CD", "trunk-based, ~32 αναπτύξεις/ημέρα"),
        ("Vertical ownership", "κάθε μηχανικός: από τη βάση δεδομένων έως το UI"),
        ("Feedback loops", "γρήγορη ανατροφοδότηση σε κάθε επίπεδο"),
    ]
    items = [{"t": "Πρακτικές που έκαναν εφικτό τον κύκλο 6 μηνών", "size": 15.5,
              "color": UP_RED, "bold": True, "font": HEAD, "space_after": 12,
              "line_spacing": 1.0}]
    for t, d in practices:
        items.append([{"t": t + ":  ", "size": 14.5, "color": UP_RED_D,
                       "bold": True, "font": BODY},
                      {"t": d, "size": 13.5, "color": DARK, "font": BODY}])
        items[-1][0]["bullet"] = True
        items[-1][0]["space_after"] = 11
        items[-1][0]["line_spacing"] = 1.05
    textbox(s, px, by + Inches(0.2), Inches(6.6), Inches(4.6), items)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE, MY CONTRIBUTIONS
# ══════════════════════════════════════════════════════════════════════════
def s_contributions():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Η συμβολή μου", "Προσωπική συμβολή στο έργο")
    bx, by, bw, bh = body_area()
    textbox(s, bx, by, Inches(12.2), Inches(0.5), [
        {"t": "Ως ιδρυτικό μέλος και τεχνικός επικεφαλής της ομάδας, συνεισέφερα "
              "καθοριστικά σε τρεις άξονες:", "size": 14.5, "color": DARK,
         "font": BODY}])
    cols = [
        ("Τεχνική θεμελίωση", [
            "Έθεσα τα θεμέλια του κώδικα και του pipeline CI/CD από την αρχή",
            "Αυτοματοποίησα το στήσιμο του περιβάλλοντος (Nix, shell): onboarding σε <30 λεπτά, χωρίς χειροκίνητα βήματα",
            "Επιτάχυνα κρίσιμα ερωτήματα με κατάλληλο indexing (GIN, trigrams) και materialized views",
        ]),
        ("Ηγεσία & ομάδα", [
            "Καθοδήγησα την ομάδα και αναπλήρωσα τον manager όποτε χρειάστηκε",
            "Σχεδίασα τη διαδικασία προσλήψεων και διεξήγαγα τις περισσότερες συνεντεύξεις",
            "Μετέδωσα τεχνογνωσία μέσα από workshops και coding dojos (event sourcing, λειτουργικός προγραμματισμός, debugging σε Elixir)",
        ]),
        ("Παράδοση & ποιότητα", [
            "Μείωσα τους χρόνους παράδοσης απλοποιώντας λύσεις και ιεραρχώντας τις ουσιώδεις απαιτήσεις",
            "Ενίσχυσα τη σαφήνεια των user stories: κριτήρια αποδοχής, προϋποθέσεις, παραδείγματα για ακραίες περιπτώσεις",
            "Ηγήθηκα της υλοποίησης του UI (Maersk Design System), σε στενή συνεργασία με τους designers",
        ]),
    ]
    cw, ch = Inches(3.92), Inches(4.55)
    gap = Inches(0.23)
    yy = by + Inches(0.62)
    for i, (head, items) in enumerate(cols):
        x = bx + i * (cw + gap)
        rect(s, x, yy, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, yy, cw, Inches(0.62), fill=UP_RED,
             shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
        textbox(s, x + Inches(0.2), yy, cw - Inches(0.4), Inches(0.62),
                [{"t": head, "size": 15.5, "color": WHITE, "bold": True,
                  "font": HEAD, "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        lines = []
        for j, it in enumerate(items):
            lines.append({"t": it, "bullet": True, "size": 12.8, "color": DARK,
                          "font": BODY, "line_spacing": 1.08,
                          "space_after": 12 if j < len(items) - 1 else 0})
        textbox(s, x + Inches(0.28), yy + Inches(0.84), cw - Inches(0.5),
                ch - Inches(1.0), lines)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 14, RESULTS: BIG NUMBERS
# ══════════════════════════════════════════════════════════════════════════
def s_results_big():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=UP_RED)
    rect(s, 0, Inches(1.12), SW, Inches(0.07), fill=UP_GOLD)
    textbox(s, Inches(0.55), Inches(0.13), Inches(11), Inches(0.3),
            [{"t": "ΑΠΟΤΕΛΕΣΜΑΤΑ", "size": 11, "color": UP_GOLD, "bold": True,
              "font": BODY}])
    textbox(s, Inches(0.55), Inches(0.40), Inches(12.2), Inches(0.66),
            [{"t": "Μετρήσιμος αντίκτυπος", "size": 25, "color": WHITE, "bold": True,
              "font": HEAD}], anchor=MSO_ANCHOR.MIDDLE)
    cards = [
        ("1.300+", "ενεργοί χρήστες", "από <100 στην εποχή Excel"),
        ("$31 εκ.", "ECO Delivery 2024", "98.000 FFE εμπορευμάτων"),
        ("6 μήνες", "έως παραγωγή", "Μάρτιος → Σεπτέμβριος 2023"),
        ("~32", "αναπτύξεις / ημέρα", "μία κάθε ~15 λεπτά"),
        ("720+", "schema migrations", "χωρίς διακοπή υπηρεσίας"),
        ("100%", "ευρημάτων GIA", "πλήρης αντιμετώπιση"),
    ]
    cw, ch = Inches(3.92), Inches(2.45)
    gx, gy = Inches(0.22), Inches(0.25)
    x0, y0 = Inches(0.55), Inches(1.55)
    for i, (big, lab, sub) in enumerate(cards):
        col, row = i % 3, i // 3
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        card = rect(s, x, y, cw, ch, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_shadow(card)
        rect(s, x, y, cw, Inches(0.12), fill=UP_GOLD)
        textbox(s, x + Inches(0.25), y + Inches(0.28), cw - Inches(0.5),
                ch - Inches(0.4), [
            {"t": big, "size": 44, "color": UP_RED, "bold": True, "font": HEAD,
             "align": PP_ALIGN.CENTER, "space_after": 2},
            {"t": lab, "size": 16, "color": DARK, "bold": True, "font": BODY,
             "align": PP_ALIGN.CENTER, "space_after": 4},
            {"t": sub, "size": 12, "color": GRAY, "italic": True, "font": BODY,
             "align": PP_ALIGN.CENTER},
        ], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, dark=True)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 15, RESULTS: FULL METRICS TABLE
# ══════════════════════════════════════════════════════════════════════════
def s_results_table():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Αποτελέσματα", "Πλήρης πίνακας μετρικών")
    bx, by, bw, bh = body_area()
    rows = [
        ["Διάσταση", "Μετρική", "Τιμή"],
        ["Υιοθέτηση", "Ενεργοί χρήστες", "1.300+ (από <100)"],
        ["Εμπορικός αντίκτυπος", "ECO Delivery Ocean 2024", "98.000 FFE / $31 εκ."],
        ["Κανονιστική συμμόρφωση", "Έλεγχος GIA 2023", "Πλήρης αντιμετώπιση"],
        ["Ταχύτητα παράδοσης", "Αναπτύξεις σε παραγωγή", "~32 / ημέρα"],
        ["Αξιοπιστία", "Schema migrations χωρίς διακοπή", "720+"],
        ["Έλεγχος ποιότητας", "Κάλυψη αυτοματοποιημένων tests", "80%+"],
        ["Κλίμακα κώδικα", "Αρχεία πηγαίου κώδικα", "~6.000"],
        ["Κύκλος ανάπτυξης", "Από έναρξη έως παραγωγή", "6 μήνες"],
    ]
    table(s, bx, by + Inches(0.1), Inches(8.1), rows, [1.15, 1.6, 1.1],
          font_size=13, row_h=Inches(0.5))
    px = Inches(9.0)
    rect(s, px, by + Inches(0.1), Inches(3.8), Inches(4.65), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, px, by + Inches(0.1), Inches(0.12), Inches(4.65), fill=UP_GOLD)
    textbox(s, px + Inches(0.34), by + Inches(0.4), Inches(3.3), Inches(4.2), [
        {"t": "Με μια ματιά", "size": 16, "color": UP_RED, "bold": True,
         "font": HEAD, "space_after": 12},
        {"t": "Από proof-of-concept σε production με μετρήσιμη εμπορική αξία σε "
              "μήνες.", "bullet": True, "size": 13.5, "color": DARK, "font": BODY,
         "space_after": 10, "line_spacing": 1.12},
        {"t": "Ταχύτητα και ποιότητα δεν είναι αντίθετοι στόχοι.", "bullet": True,
         "size": 13.5, "color": DARK, "font": BODY, "space_after": 10,
         "line_spacing": 1.12},
        {"t": "Ο τεχνικός σχεδιασμός υπηρετεί και την εταιρική διακυβέρνηση.",
         "bullet": True, "size": 13.5, "color": DARK, "font": BODY,
         "line_spacing": 1.12},
    ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 16, LESSONS
# ══════════════════════════════════════════════════════════════════════════
def s_lessons():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Συμπεράσματα", "Μαθήματα & συμβιβασμοί (trade-offs)")
    bx, by, bw, bh = body_area()
    lessons = [
        ("Vertical ownership vs γνωσιακό φορτίο",
         "Μεγάλη συνοχή, αλλά απαιτεί ευρεία γνώση, αντίμετρο το pairing & onboarding"),
        ("Monorepo: ισχύς με κόστος",
         "Ατομικές αναπτύξεις & refactoring, αλλά αυξανόμενος χρόνος build/CI"),
        ("Event sourcing: επιλεκτικά",
         "Ιδανικό για audit-heavy προβλήματα (Energy Bank), αλλά αυξάνει την καμπύλη μάθησης"),
        ("Cloud vs self-hosting",
         "Ευελιξία κλιμάκωσης, αλλά κόστος εξαρτώμενο από τον όγκο, απαιτεί σχεδιασμό"),
        ("Ταχύτητα ως πλεονέκτημα",
         "TDD + pair programming + CD: η ποιότητα επιτρέπει την ταχύτητα, δεν την εμποδίζει"),
    ]
    yy = by + Inches(0.15)
    for i, (t, d) in enumerate(lessons):
        rect(s, bx, yy, Inches(12.2), Inches(0.82), fill=LIGHT if i % 2 == 0 else WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, bx, yy, Inches(0.1), Inches(0.82), fill=UP_RED)
        textbox(s, bx + Inches(0.3), yy, Inches(11.7), Inches(0.82), [
            [{"t": t + "    ", "size": 15, "color": UP_RED_D, "bold": True,
              "font": HEAD},
             {"t": d, "size": 13, "color": DARK, "font": BODY}]],
            anchor=MSO_ANCHOR.MIDDLE)
        yy = yy + Inches(0.92)
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 17, FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════
def s_future():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=WHITE)
    header(s, "Μελλοντική εργασία", "Επόμενα βήματα")
    bx, by, bw, bh = body_area()
    future = [
        ("Scope 2 εκπομπές", "Γραφεία, τερματικά και cold ironing πλοίων"),
        ("Επέκταση κάλυψης στόλου", "Περισσότερες δραστηριότητες & τύποι μεταφοράς"),
        ("Επέκταση αναφοράς SBTi", "Πληρέστερη παρακολούθηση κλιματικών δεσμεύσεων"),
        ("Προβλεπτικά analytics", "Σενάρια αποκαρβονοποίησης & πρόβλεψη εκπομπών"),
        ("Bunker ↔ Energy Bank", "Σύνδεση βελτιστοποίησης με τη λογιστική πράσινου καυσίμου"),
    ]
    cw, ch = Inches(3.92), Inches(2.0)
    gx, gy = Inches(0.22), Inches(0.25)
    x0, y0 = bx, by + Inches(0.3)
    for i, (t, d) in enumerate(future):
        col, row = i % 3, i // 3
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, x + Inches(0.28), y + Inches(0.22), cw - Inches(0.5),
                ch - Inches(0.4), [
            {"t": "→", "size": 24, "color": UP_GOLD, "bold": True, "font": HEAD,
             "space_after": 4},
            {"t": t, "size": 15.5, "color": UP_RED_D, "bold": True, "font": HEAD,
             "space_after": 6, "line_spacing": 1.0},
            {"t": d, "size": 12.5, "color": DARK, "font": BODY, "line_spacing": 1.08},
        ])
    footer(s)


# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 18, CLOSING
# ══════════════════════════════════════════════════════════════════════════
def s_closing():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=LIGHT)
    rect(s, 0, 0, Inches(0.32), SH, fill=UP_RED)
    rect(s, Inches(0.32), 0, Inches(0.07), SH, fill=UP_GOLD)
    # University emblem + department logo, side by side at the top
    add_image(s, os.path.join(ASSETS, "logo-up-emblem.png"),
              Inches(3.76), Inches(0.62), Inches(1.5), Inches(1.5), center_in=False)
    rect(s, Inches(5.61), Inches(0.85), Pt(2), Inches(1.04), fill=UP_GOLD)
    add_image(s, os.path.join(ASSETS, "logo-ceid.png"),
              Inches(5.98), Inches(0.94), Inches(3.6), Inches(0.86), center_in=True)
    textbox(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(1.2), [
        {"t": "Σας ευχαριστώ", "size": 40, "color": UP_RED_D, "bold": True,
         "font": HEAD, "align": PP_ALIGN.CENTER, "space_after": 6},
        {"t": "Ερωτήσεις & συζήτηση", "size": 20, "color": GRAY, "italic": True,
         "font": BODY, "align": PP_ALIGN.CENTER}])
    rect(s, Inches(4.6), Inches(4.95), Inches(4.1), Pt(2), fill=UP_GOLD)
    textbox(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.4), [
        {"t": "Κωνσταντίνος Γκίνης", "size": 17, "color": DARK, "bold": True,
         "font": HEAD, "align": PP_ALIGN.CENTER, "space_after": 4},
        {"t": "Επιβλέπων: Χρήστος Μπούρας, Καθηγητής", "size": 13, "color": GRAY,
         "font": BODY, "align": PP_ALIGN.CENTER, "space_after": 2},
        {"t": "Πανεπιστήμιο Πατρών · Τμήμα Μηχανικών Η/Υ & Πληροφορικής · Ιούνιος 2026",
         "size": 12, "color": GRAY, "font": BODY, "align": PP_ALIGN.CENTER}])


# ══════════════════════════════════════════════════════════════════════════
#  APPENDIX, full-screen versions of the screenshots (click-to-enlarge)
# ══════════════════════════════════════════════════════════════════════════
def build_zoom_appendix():
    """For every registered screenshot, add a full-screen slide after the deck
    and link thumbnail <-> full version both ways. The appendix sits past the
    closing slide, so it never interrupts the linear flow."""
    for z in ZOOMS:
        ap = slide()
        bg = rect(ap, 0, 0, SW, SH, fill=RGBColor(0x1E, 0x1E, 0x22))
        big = add_image(ap, z["path"], Inches(0.3), Inches(0.18),
                        Inches(12.73), Inches(6.5), center_in=True)
        textbox(ap, Inches(0.45), Inches(6.96), Inches(8.5), Inches(0.42),
                [{"t": z["caption"], "size": 12.5, "color": WHITE, "bold": True,
                  "font": BODY}], anchor=MSO_ANCHOR.MIDDLE)
        textbox(ap, Inches(8.5), Inches(6.96), Inches(4.4), Inches(0.42),
                [{"t": "κλικ οπουδήποτε για επιστροφή", "size": 11,
                  "color": UP_GOLD, "italic": True, "font": BODY,
                  "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
        # thumbnail -> full screen
        z["pic"].click_action.target_slide = ap
        if z["bar"] is not None:
            z["bar"].click_action.target_slide = ap
        # full screen -> back to the source slide (any click returns)
        bg.click_action.target_slide = z["src"]
        big.click_action.target_slide = z["src"]


def add_card_shadows():
    """Give every beige card/panel the same soft drop shadow as the screenshot
    frames. The width filter skips full-width callout bars, the zebra-striped
    lessons rows and the full-slide backgrounds."""
    beige = (str(LIGHT), str(PANEL))
    for sl in prs.slides:
        for sp in sl.shapes:
            if sp.width >= Inches(9.5):
                continue
            try:
                rgb = str(sp.fill.fore_color.rgb)
            except Exception:
                continue
            if rgb in beige:
                _add_shadow(sp)


# ── build ─────────────────────────────────────────────────────────────────
s_title()
s_problem()
s_audit()
s_func_req()
s_nonfunc_req()
s_pillars()
s_arch()
s_demo1()
s_demo2()
s_demo3()
s_theory()
s_tech()
s_code()
s_methods()
s_contributions()
s_results_big()
s_lessons()
s_future()
s_closing()
build_zoom_appendix()
add_card_shadows()


# ── Speaker notes (Greek), budgeted to ~19 min + buffer for a 20΄ defense ──
NOTES = [
    # 1, Title
    """⏱ ~30΄΄  |  Σύνολο στόχος: 20 λεπτά.
Καλησπέρα σας. Ονομάζομαι Κωνσταντίνος Γκίνης. Η διπλωματική μου παρουσιάζει
τον «Πάγκο Εργασίας Εκπομπών» (Emissions Workbench): μια πλατφόρμα λογισμικού
παραγωγικής κλίμακας που σχεδίασε και υλοποίησε η ομάδα Energy Transition
Platform της Maersk, για τη μέτρηση, πιστοποίηση και βελτιστοποίηση των εκπομπών
αερίων θερμοκηπίου στη ναυτιλία. Επιβλέπων ο καθηγητής κ. Χρήστος Μπούρας.
→ Ξεκινώ από το πρόβλημα που γέννησε την ανάγκη.""",
    # 2, Problem
    """⏱ ~70΄΄
Η Maersk, με στόλο 700+ πλοίων, δεν μπορούσε να ανταποκριθεί αξιόπιστα στις
αυξανόμενες απαιτήσεις για δεδομένα εκπομπών. Η διαδικασία ήταν εξ ολοκλήρου
χειρωνακτική: αίτημα μέσω email, μη αυτόματη εξαγωγή δεδομένων, εφαρμογή
συντελεστών σε φύλλα Excel, και απάντηση μετά από 3 έως 8 εβδομάδες.
Αυτό κόστιζε σε τέσσερις άξονες: ΑΚΡΙΒΕΙΑ (λάθη μεθοδολογίας), ΤΑΧΥΤΗΤΑ
(καθυστερήσεις που σήμαιναν απώλεια συμβολαίων), ΙΧΝΗΛΑΣΙΜΟΤΗΤΑ (κανένα
audit) και ΚΛΙΜΑΚΑ (η ανθρώπινη δυναμικότητα δεν κλιμακώνεται για 1.300+ χρήστες).
→ Το πρόβλημα ήταν γνωστό· χρειαζόταν όμως ένας καταλύτης.""",
    # 3, Audit
    """⏱ ~70΄΄
Αυτός ο καταλύτης ήταν ο εσωτερικός έλεγχος Group Internal Audit του 2023.
Εντόπισε τέσσερα συγκεκριμένα ευρήματα: έλλειψη παρακολούθησης δεδομένων,
απουσία ίχνους ελέγχου (κρίσιμο εν όψει CSRD και EU ETS), αδυναμία κλιμάκωσης,
και υψηλό ποσοστό σφαλμάτων. Η σημασία του ελέγχου ήταν ότι έδωσε επίσημο
βάρος σε ένα ήδη γνωστό πρόβλημα, και έτσι εγκρίθηκαν πόροι για ένα εξ
ολοκλήρου νέο σύστημα. Η ανάπτυξη ξεκίνησε τον Μάρτιο 2023 και έφτασε σε
παραγωγή τον Σεπτέμβριο 2023: μόλις έξι μήνες.
→ Ας δούμε τι έπρεπε να κάνει αυτό το σύστημα, τις προδιαγραφές.""",
    # 4, Functional requirements
    """⏱ ~75΄΄
Οι λειτουργικές απαιτήσεις προέκυψαν από έξι διακριτές ομάδες χρηστών, η καθεμία
με διαφορετική ανάγκη και συχνότητα: από τα ECO Delivery Products που θέλουν
αυτόματα πιστοποιητικά πολλές φορές την ημέρα, μέχρι τους Account Managers που
χρειάζονται δεδομένα σε πραγματικό χρόνο μέσα σε μια διαπραγμάτευση.
Από αυτή την ποικιλομορφία προκύπτουν οι βασικές λειτουργίες του συστήματος:
αναζήτηση εκπομπών ανά πελάτη και χρονικό εύρος, ανάλυση ανά αποστολή και
δρομολόγιο, εξαγωγή σε Excel, αυτόματη έκδοση πιστοποιητικών, και όλα αυτά
σε λογική self-service, χωρίς μεσολάβηση ειδικής ομάδας αναφορών.
→ Πέρα όμως από το «τι κάνει», υπάρχει το «πώς πρέπει να συμπεριφέρεται».""",
    # 5, Non-functional requirements
    """⏱ ~75΄΄
Εδώ είναι το σημείο που ζήτησε ιδιαίτερα η επιτροπή: οι μη-λειτουργικές, τεχνικές
απαιτήσεις. Συνδυάζοντας τις ανάγκες των χρηστών με τα ευρήματα GIA, το σύστημα
έπρεπε να εγγυάται: πλήρη ιχνηλασιμότητα δεδομένων· αμετάβλητο audit trail·
δεδομένα πραγματικού χρόνου (freshness σε ώρες, όχι εβδομάδες)· κλιμακωσιμότητα·
υψηλή διαθεσιμότητα· επαληθευσιμότητα για CSRD, EU ETS, ISCC· πλήρη
αυτοματοποίηση· και ασφαλή εξουσιοδότηση.
Καθεμία από αυτές καθόρισε μια αρχιτεκτονική επιλογή, που φαίνεται στην κάτω
γραμμή: event-driven αρχιτεκτονική, event sourcing για το audit trail, Phoenix
LiveView για real-time, και BEAM/Elixir για διαθεσιμότητα.
→ Δείτε πώς υλοποιήθηκαν αυτές οι απαιτήσεις, η λύση.""",
    # 6, Pillars
    """⏱ ~80΄΄
Η λύση είναι μία ενιαία πλατφόρμα, σε ένα monorepo, οργανωμένη σε τέσσερις
πυλώνες. Ocean Emissions: αυτόματη μέτρηση εκπομπών ανά αποστολή, από τηλεμετρία
στόλου (STAR Connect) και ενσωμάτωση EU ETS. ECO Delivery: τα εμπορικά προϊόντα
χαμηλού άνθρακα, με αυτόματα πιστοποιητικά και τη λογιστική πράσινου καυσίμου του
Energy Bank. Net Zero 2040: εσωτερικοί πίνακες για οδικούς χάρτες αποκαρβονοποίησης.
Bunker Optimization: βελτιστοποίηση πλάνων ανεφοδιασμού για ελάχιστο κόστος καυσίμου.
Σημαντικό: οι πυλώνες δεν είναι ανεξάρτητοι, επικοινωνούν μέσω Kafka, όπως
δείχνει η ροή δεδομένων στο κάτω μέρος.
→ Πώς οργανώνεται τεχνικά όλο αυτό;""",
    # 7, Architecture
    """⏱ ~70΄΄
Η αρχιτεκτονική στηρίζεται σε πέντε αρχές. Clean Architecture: ο επιχειρησιακός
πυρήνας είναι ανεξάρτητος από την υποδομή. Event-driven ingestion μέσω Kafka και
εργασιών Oban. Event Sourcing ειδικά για το Energy Bank, που δίνει πλήρες ιστορικό
και το αμετάβλητο audit trail. Real-time web με Phoenix LiveView, χωρίς reload
σελίδας. Και ανάπτυξη σε Azure/Kubernetes με auto-scaling και infrastructure-as-code
μέσω Terraform. Το διάγραμμα δείχνει τη συνολική ροή, από τις πηγές δεδομένων,
μέσα από τα components, έως το επίπεδο ιστού.
→ Ας δούμε τώρα την πλατφόρμα στην πράξη.""",
    # 8, Demo 1
    """⏱ ~70΄΄  (test data)
Η πλατφόρμα στην πράξη. Αριστερά, η οθόνη αναζήτησης: ο χρήστης εισάγει κωδικό
πελάτη και χρονικό εύρος, και τα συγκεντρωτικά δεδομένα εκπομπών (όγκος, εκπομπές,
εξοικονόμηση) επιστρέφονται άμεσα, αντικαθιστώντας τη ροή email-και-αναμονής 3 έως
8 εβδομάδων. Τεχνικά, γίνεται polling από το Dremio μέσω εργασιών Oban, που
εξασφαλίζει freshness σε ώρες. Δεξιά, το drill-down ανά αποστολή και δρομολόγιο:
ο πελάτης βλέπει ποια trade lanes βαραίνουν περισσότερο το αποτύπωμά του, που
είναι η βάση για να προτείνουμε υπηρεσίες ECO.
→ Από τα δεδομένα, περνάμε στα πιστοποιητικά.""",
    # 9, Demo 2
    """⏱ ~60΄΄
Αριστερά, το παραγόμενο πιστοποιητικό ECO σε PDF, με αναλυτική αναφορά
μεθοδολογίας ώστε ο πελάτης να το χρησιμοποιήσει στις δικές του αναφορές CSRD.
Δεξιά πάνω, ο πίνακας του Oban: οι ασύγχρονες εργασίες που τρέχουν στο παρασκήνιο,
τόσο η άντληση δεδομένων (ingestion από Dremio) όσο και η παραγωγή πιστοποιητικών,
οργανωμένες σε ουρές. Αυτές οι εργασίες αντλούν από το ισοζύγιο του Energy Bank,
επαληθεύουν τα διαθέσιμα υπόλοιπα βιώσιμου καυσίμου, και παράγουν έγγραφα ISCC.
Δεξιά κάτω, η εξαγωγή σε Excel, τη μορφή που χρησιμοποιούν στην πράξη οι τελικοί
χρήστες. Όλη αυτή η ροή, που πριν απαιτούσε χειρωνακτική εργασία εβδομάδων,
εκτελείται πλέον αυτόματα, σε λεπτά.
→ Πίσω από αυτά υπάρχουν ακόμη δύο σημαντικά υποσυστήματα.""",
    # 10, Demo 3
    """⏱ ~50΄΄
Δύο ακόμη όψεις πίσω από τη διεπαφή. Το Energy Bank: ένα event-sourced ledger με
καταθέσεις και αναλήψεις βιώσιμου καυσίμου, που επιβάλλει τον περιορισμό του
ισοζυγίου μάζας. Και το Bunker Optimization: το αποτέλεσμα του solver, δηλαδή το
βέλτιστο πλάνο ανεφοδιασμού ανά λιμάνι και δρομολόγιο.
→ Πριν τα αποτελέσματα, μια σύντομη θεωρητική βάση.""",
    # 11, Theory
    """⏱ ~60΄΄  [σύντομα· η επιτροπή ζήτησε λιγότερη θεωρία]
Πολύ συνοπτικά το θεωρητικό υπόβαθρο. Οι εκπομπές κατηγοριοποιούνται σε Scope 1,
2 και 3 (άμεσες, έμμεσες από ενέργεια, αλυσίδας αξίας). Τα όρια του συστήματος
ορίζονται ως Tank-to-Wake ή Well-to-Wake. Ο πυρήνας του υπολογισμού είναι οι
trade factors, οι συντελεστές εμπορικής διαδρομής. Και η λογιστική πράσινου
καυσίμου στηρίζεται στο mass balance και στο book-and-claim, υπό τα πρότυπα
ISCC, GLEC και το GHG Protocol.
Το ουσιώδες: αυτές οι έννοιες δεν έμειναν θεωρητικές, αντιστοιχήθηκαν άμεσα σε
σχεδιαστικά πρότυπα μέσα στον κώδικα.
→ Ποια τεχνολογία τα υλοποιεί;""",
    # 12, Tech
    """⏱ ~60΄΄
Η στοίβα στηρίζεται σε τρεις πυλώνες: την πλατφόρμα BEAM/Elixir για διαθεσιμότητα
και ταυτοχρονισμό· το Phoenix LiveView για real-time διεπαφή χωρίς reload· και το
Nix για αναπαραγώγιμα builds και περιβάλλοντα ανάπτυξης. Τα υποστηρικτικά:
PostgreSQL/Ecto για μονιμότητα δεδομένων, Apache Kafka για την event-driven ροή
μηνυμάτων, GitHub Actions για CI/CD, και Azure με Kubernetes και Terraform για
deployment και infrastructure-as-code.
→ Εξίσου σημαντικό με την τεχνολογία ήταν ο τρόπος εργασίας της ομάδας.""",
    # 12b, Code highlights
    """⏱ ~55΄΄  [προαιρετική: παράλειψε αν πιέζει ο χρόνος]
Δύο σύντομα παραδείγματα από τον πραγματικό κώδικα, που δείχνουν γιατί ο
λειτουργικός προγραμματισμός ταιριάζει στο πρόβλημα. Πάνω αριστερά: το pattern
matching σε πολλαπλές ρήτρες της calculate επιλέγει αυτόματα τη σωστή περίπτωση
και κανονικοποιεί τις μονάδες (π.χ. megajoules σε kWh) δηλωτικά, χωρίς if. Κάτω
αριστερά: ο τελεστής pipe διοχετεύει μια τιμή μέσα από διαδοχικούς
μετασχηματισμούς, αριστερά προς δεξιά, και ένα ερώτημα Ecto που είναι δηλωτικό
και type-safe. Δεξιά συνοψίζονται τα οφέλη, με κυριότερο την αμεταβλητότητα, που
δίνει ασφάλεια στον ταυτοχρονισμό του BEAM.
→ Περνάω στις εργασιακές μεθόδους.""",
    # 13, Methods
    """⏱ ~70΄΄
Η ομάδα των 28 μελών εφάρμοσε Extreme Programming. Pair programming: συνεχής
έλεγχος κώδικα και μεταφορά γνώσης. TDD, με κάλυψη 80%+. Continuous Integration/
Delivery σε λογική trunk-based, που έφτανε τις ~32 αναπτύξεις την ημέρα. Vertical
ownership: κάθε μηχανικός αναλαμβάνει λειτουργικότητα από τη βάση δεδομένων έως το
UI. Και σφιχτοί βρόχοι ανατροφοδότησης σε κάθε επίπεδο. Αυτές οι πρακτικές είναι
που έκαναν εφικτό τον κύκλο των έξι μηνών έως την παραγωγή.
→ Μέσα σε αυτό το πλαίσιο, ποια ήταν η δική μου συμβολή;""",
    # 15, My contributions
    """⏱ ~90΄΄  [η προσωπική διαφάνεια· μίλησε με σιγουριά και σε πρώτο πρόσωπο]
Ως ιδρυτικό μέλος και τεχνικός επικεφαλής, η συμβολή μου εκτείνεται σε τρεις άξονες.
ΤΕΧΝΙΚΗ ΘΕΜΕΛΙΩΣΗ: ξεκίνησα τον κώδικα και έστησα το CI/CD pipeline· έγραψα tooling
με Nix και shell που στήνει το περιβάλλον ανάπτυξης ιδιοτελώς (idempotent),
φέρνοντας το onboarding σε λιγότερο από 30 λεπτά χωρίς καμία χειρωνακτική παρέμβαση·
και βελτιστοποίησα κρίσιμα ερωτήματα με καλύτερο indexing (GIN, trigrams) και
materialized views. ΗΓΕΣΙΑ ΚΑΙ ΟΜΑΔΑ: καθοδήγησα την ομάδα και αναπλήρωσα τον
manager όταν χρειάστηκε· σχεδίασα τη διαδικασία προσλήψεων και έκανα τις
περισσότερες συνεντεύξεις· και οργάνωσα workshops και coding dojos για event
sourcing, λειτουργικό προγραμματισμό και debugging σε Elixir. ΠΑΡΑΔΟΣΗ ΚΑΙ
ΠΟΙΟΤΗΤΑ: μείωσα τους χρόνους παράδοσης απλουστεύοντας λύσεις και δίνοντας
προτεραιότητα στις σημαντικότερες απαιτήσεις· βελτίωσα τη σαφήνεια των user stories
με ρητά acceptance criteria, προϋποθέσεις και παραδείγματα για edge cases· και
ηγήθηκα της υλοποίησης των UI components του Maersk Design System σε συνεργασία με
τους designers.
→ Και ερχόμαστε στο πιο σημαντικό: τα αποτελέσματα.""",
    # 16, Results (big numbers)
    """⏱ ~80΄΄  [η κορυφαία διαφάνεια· δώσε χρόνο και έμφαση]
Ο αντίκτυπος είναι μετρήσιμος. 1.300+ ενεργοί χρήστες, από λιγότερους από 100
στην εποχή του Excel. ECO Delivery αξίας 31 εκατομμυρίων δολαρίων, ή 98.000 FFE,
μόνο για το 2024, μια εμπορική δραστηριότητα που πριν δεν ήταν τεχνικά εφικτή σε
αυτή την κλίμακα. Από την έναρξη έως την παραγωγή: έξι μήνες. ~32 αναπτύξεις την
ημέρα, μία κάθε δεκαπέντε λεπτά. 720+ schema migrations χωρίς καμία διακοπή
υπηρεσίας. Και πλήρης, 100%, αντιμετώπιση όλων των ευρημάτων του ελέγχου GIA.
Πέρα από τους αριθμούς, δύο συμπεράσματα: ταχύτητα και ποιότητα δεν είναι αντίθετοι
στόχοι, και ο τεχνικός σχεδιασμός (event sourcing, audit trail) υπηρετεί και την
εταιρική διακυβέρνηση.
→ Τι μάθαμε από την υλοποίηση;""",
    # 17, Lessons
    """⏱ ~70΄΄
Πέντε μαθήματα, κυρίως συμβιβασμοί. Το vertical ownership δίνει συνοχή, αλλά αυξάνει
το γνωσιακό φορτίο, αντίμετρο το pairing και το onboarding. Το monorepo προσφέρει
ατομικές αναπτύξεις, με κόστος χρόνου build. Το event sourcing είναι ιδανικό για
audit-heavy προβλήματα όπως το Energy Bank, αλλά πρέπει να εφαρμόζεται επιλεκτικά.
Το cloud δίνει ευελιξία, με κόστος που απαιτεί σχεδιασμό. Και το βασικότερο: η
ταχύτητα παράδοσης, όταν στηρίζεται σε TDD και pair programming, γίνεται
ανταγωνιστικό πλεονέκτημα, η ποιότητα επιτρέπει την ταχύτητα, δεν την εμποδίζει.
→ Πού πηγαίνει το έργο από εδώ;""",
    # 17, Future work
    """⏱ ~50΄΄
Οι κυριότερες επεκτάσεις που ήδη σχεδιάζονται: κάλυψη εκπομπών Scope 2 (γραφεία,
τερματικά, cold ironing)· επέκταση κάλυψης στόλου και δραστηριοτήτων· πληρέστερη
αναφορά SBTi· προβλεπτικά analytics με σενάρια αποκαρβονοποίησης· και στενότερη
σύνδεση του Bunker Optimization με τη λογιστική του Energy Bank.
→ Κλείνω.""",
    # 18, Closing
    """⏱ ~30΄΄
Συνοψίζοντας: ο Πάγκος Εργασίας Εκπομπών αποδεικνύει ότι ένα σύστημα μπορεί να
συνδυάσει επιστημονική ακρίβεια, εμπορική αξιοπιστία και υψηλό ρυθμό παράδοσης ,
αρκεί οι τεχνολογικές, αρχιτεκτονικές και μεθοδολογικές επιλογές να αντιμετωπίζονται
ως ενιαίο σύστημα. Σας ευχαριστώ για την προσοχή σας· είμαι στη διάθεσή σας για
ερωτήσεις.""",
]

for _slide, _note in zip(prs.slides, NOTES):
    _slide.notes_slide.notes_text_frame.text = _note

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst),
      "| notes:", len(NOTES))
